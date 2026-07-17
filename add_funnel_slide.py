#!/usr/bin/env python3
# -*- coding: utf-8 -*-
"""Add a data-funnel slide (2,068 raw deals -> 73 usable labels) and place it before the
updated-pipeline slide. Reuses update_deck.py's design-token helpers."""
from pptx import Presentation
from pptx.enum.text import PP_ALIGN
import update_deck as ud

SRC = "JPM_Project_3ii_week6.pptx"


def title_of(slide):
    for sh in slide.shapes:
        if sh.has_text_frame and sh.text_frame.text.strip():
            return sh.text_frame.text
    return ""


def main():
    prs = Presentation(SRC)
    s = ud.blank_slide(prs)
    ud.scaffold(s, "From 2,068 raw deals to 73 usable labels",
                "Most of the funnel is structural — deals that were never eligible, not pipeline losses")

    rows = [
        ["Filter", "Remaining", "Dropped", "Why it drops"],
        ["Raw Bloomberg pull (2006+)", "2,068", "—", "Everything the M&A screen returned"],
        ["Keep “Cash or Stock” (true election)", "727", "−1,341", "Fixed-mix deals — no election, no proration"],
        ["Keep completed", "619", "−108", "Terminated / withdrawn — no election happened"],
        ["US + resolvable identifier", "317", "−302", "Non-US — no EDGAR filings to read"],
        ["Ran through EDGAR + Claude", "303", "−14", "No CIK / no election filing found"],
        ["Clean disclosed election demand", "73", "−231", "Not disclosed, or unparseable prose"],
    ]
    ud.table(s, 0.70, 2.02, 11.95, rows, [4.30, 1.30, 1.30, 5.05], row_h=0.44)

    # takeaway callout
    ud.rect(s, 0.70, 5.78, 11.95, 0.92, ud.BOXBG)
    ud.txt(s, "Where the deals go", 0.86, 5.90, 11.6, 0.20, 9.5, ud.NAVY, bold=True)
    ud.txt(s,
           "The two largest cuts are structural, not pipeline losses — 1,341 “Cash-and-Stock” deals have no "
           "election at all, and 302 non-US deals have no EDGAR filings (together ~79% of the drop, right at the top). "
           "The remaining 73 are near the disclosure ceiling: election demand simply isn’t disclosed on most other "
           "deals — a sharp-prompt re-extraction of 26 candidates recovered only 1.",
           0.86, 6.14, 11.63, 0.50, 9.6, ud.DARK)

    # page-number placeholder (renumber() will set the real value)
    ud.txt(s, "0 / 0", 12.10, 7.08, 0.80, 0.20, 7.8, ud.SLATE, align=PP_ALIGN.RIGHT)

    # move the new (last) slide to sit just before the "Updated pipeline" slide
    lst = prs.slides._sldIdLst
    new = list(lst)[-1]
    pos = next(i for i, sl in enumerate(prs.slides) if "Updated pipeline" in title_of(sl))
    lst.remove(new)
    lst.insert(pos, new)

    ud.renumber(prs)
    prs.save(SRC)
    print(f"[funnel] added. deck now {len(prs.slides)} slides:")
    for i, sl in enumerate(prs.slides):
        print(f"  {i:2d}  {title_of(sl).split(chr(10))[0][:52]}")


if __name__ == "__main__":
    main()
