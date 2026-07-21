#!/usr/bin/env python3
# -*- coding: utf-8 -*-
"""Update the week-6 deck: refresh the Weekly Update, add an updated-pipeline slide and a
one-deal extraction walkthrough slide. Matches the deck's existing design tokens exactly."""
import re
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE

NAVY = RGBColor(0x14, 0x31, 0x5A); NAVY2 = RGBColor(0x18, 0x2B, 0x49)
SLATE = RGBColor(0x47, 0x55, 0x69); BLUE = RGBColor(0x25, 0x63, 0xEB)
DARK = RGBColor(0x0F, 0x17, 0x2A); WHITE = RGBColor(0xFF, 0xFF, 0xFF)
BOXBG = RGBColor(0xEA, 0xF2, 0xFF); ZEBRA = RGBColor(0xF7, 0xF9, 0xFC)
BORDER = RGBColor(0xE2, 0xE8, 0xF0)

SRC = "JPM_Project_3ii_week7.pptx"


def txt(slide, text, l, t, w, h, size, color, bold=False, align=PP_ALIGN.LEFT):
    tb = slide.shapes.add_textbox(Inches(l), Inches(t), Inches(w), Inches(h))
    tf = tb.text_frame; tf.word_wrap = True
    tf.margin_left = 0; tf.margin_right = 0; tf.margin_top = 0; tf.margin_bottom = 0
    p = tf.paragraphs[0]; p.alignment = align
    r = p.add_run(); r.text = text
    r.font.size = Pt(size); r.font.bold = bold; r.font.color.rgb = color
    return tb


def rect(slide, l, t, w, h, fill, line=None):
    sh = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(l), Inches(t), Inches(w), Inches(h))
    sh.fill.solid(); sh.fill.fore_color.rgb = fill
    if line is not None:
        sh.line.color.rgb = line; sh.line.width = Pt(0.75)
    else:
        sh.line.fill.background()
    sh.shadow.inherit = False
    return sh


def blank_slide(prs):
    s = prs.slides.add_slide(prs.slide_layouts[0])
    for ph in list(s.placeholders):
        ph._element.getparent().remove(ph._element)
    return s


def scaffold(slide, title, subtitle):
    rect(slide, 0, 0, 13.33, 7.5, WHITE)
    rect(slide, 0, 0, 13.33, 0.18, NAVY)
    txt(slide, title, 0.55, 0.40, 12.20, 0.55, 28, NAVY, bold=True)
    txt(slide, subtitle, 0.56, 1.01, 11.90, 0.36, 13.5, SLATE)
    txt(slide, "JPM Project 3ii · Weekly workflow validation", 0.55, 7.08, 7.5, 0.20, 7.8, SLATE)


def box(slide, l, t, w, h, title, sub):
    rect(slide, l, t, w, h, BOXBG)
    txt(slide, title, l + 0.09, t + 0.13, w - 0.18, 0.20, 8.8, NAVY, bold=True, align=PP_ALIGN.CENTER)
    txt(slide, sub, l + 0.09, t + 0.39, w - 0.18, 0.20, 7.2, SLATE, align=PP_ALIGN.CENTER)


def table(slide, l, t, w, rows, colw, header_h=0.42, row_h=0.40, fs=8.8):
    """rows[0] = header; colw = list of column widths summing to w."""
    y = t
    for ri, row in enumerate(rows):
        fill = NAVY if ri == 0 else (WHITE if (ri % 2 == 1) else ZEBRA)
        rect(slide, l, y, w, header_h if ri == 0 else row_h, fill)
        x = l + 0.08
        for ci, cell in enumerate(row):
            col = WHITE if ri == 0 else DARK
            txt(slide, cell, x, y + (0.13 if ri == 0 else 0.11), colw[ci] - 0.12,
                0.20, 8.9 if ri == 0 else fs, col, bold=(ri == 0))
            x += colw[ci]
        y += header_h if ri == 0 else row_h
    return y


def update_weekly(prs):
    s = prs.slides[1]
    body = s.shapes[2]
    body.top = Inches(1.70); body.left = Inches(0.75); body.width = Inches(11.9)
    tf = body.text_frame; tf.clear()
    lines = [
        ("Accomplished this week", True),
        ("Completed the full-universe SEC → Claude extraction: 317 deals, 303 with fields, 73 clean election-demand labels", False),
        ("Built the Monte Carlo framework — demand distribution (Beta) → proration mechanics → simulated payoff and portfolio P&L", False),
        ("Validated the demand model out-of-sample: leave-one-out calibration passes (KS p = 0.96), model is honest, not over-confident", False),
        ("Built the trade decision layer (blotter): entry, election side, hedge, sizing and a go/no-go rule — signal skill corr +0.67", False),
        ("Recovered missing acquirer prices via CUSIP → PERMNO (identifier-drift fix, e.g. BB&T→Truist); MC-ready deals 25 → 32", False),
        ("Proved the demand-disclosure ceiling empirically — a sharp-prompt re-extraction recovered ~0; it is non-disclosure, not a tooling gap", False),
        ("In progress / to do", True),
        ("Survivorship-aware P&L backtest: add the terminated deals and extend prices to close so deal-break risk is priced, not assumed", False),
        ("Calibrate the deal-break probability and entry hurdle from data (currently placeholders: 12% break, 0.5% hurdle)", False),
        ("Censored-demand refinement ($0): fold proration-outcome deals into the Beta fit as bounded demand observations", False),
        ("Extend to the private/foreign-acquirer deals as an 'unhedgeable' bucket, and write up the two structural ceilings", False),
    ]
    for i, (text, bold) in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        r = p.add_run(); r.text = text if bold else ("•  " + text)
        r.font.size = Pt(12); r.font.bold = bold
        r.font.color.rgb = NAVY2 if bold else DARK
        p.space_after = Pt(3)
        if bold and i > 0:
            p.space_before = Pt(8)


def build_pipeline(prs):
    s = blank_slide(prs)
    scaffold(s, "Updated pipeline — extraction to Monte Carlo to trade",
             "Full current flow: 2,068 raw deals → 73 clean labels → a calibrated model → a trade blotter")
    stages = [("Universe", "2,068 → 317"), ("SEC + Claude", "field extraction"),
              ("WRDS", "prices + own."), ("Labels", "73 clean demand"),
              ("Monte Carlo", "demand → payoff"), ("Trade blotter", "entry + hedge")]
    x = 0.65
    for i, (a, b) in enumerate(stages):
        box(s, x, 1.62, 1.55, 0.74, a, b)
        if i < len(stages) - 1:
            rect(s, x + 1.60, 1.90, 0.16, 0.20, RGBColor(0xCB, 0xD5, 0xE1))
        x += 2.05
    rows = [
        ["Stage", "Primary output", "Current state"],
        ["0–1  Universe + identifiers", "US_election_deals_for_analysis.csv (317)", "CIK + acquirer-PERMNO overrides for delisted/renamed"],
        ["2  SEC + Claude extraction", "ma_edgar_full/llm_field_extractions.csv", "303 deals, 17 canonical fields each"],
        ["3  WRDS market + ownership", "wrds_market_daily.csv · ownership_mix", "Target + acquirer prices, passive block"],
        ["4  Normalize labels", "normalized_labels.csv", "73 clean election-DEMAND labels (near disclosure ceiling)"],
        ["5  Deadline spread", "deadline_spread.csv", "Fixed vs floating split (floating = spread≈0, excluded)"],
        ["6  Monte Carlo  (arb_mc / arb_run)", "arb_output/ figures + summary", "Demand Beta; calibration KS p = 0.96"],
        ["7  Trade decision  (arb_signal)", "arb_signals.csv", "31 deals, 22 ENTER, signal skill corr +0.67"],
        ["8  Deliverable", "arb_output/walkthrough.html", "This deck + the browser walkthrough artifact"],
    ]
    table(s, 0.70, 2.75, 11.95, rows, [3.30, 3.95, 4.70])
    txt(s, "22 / 14", 12.10, 7.08, 0.80, 0.20, 7.8, SLATE, align=PP_ALIGN.RIGHT)
    return s


def build_walkthrough(prs):
    s = blank_slide(prs)
    scaffold(s, "Extraction walkthrough — one deal, end to end",
             "PepsiAmericas ← PepsiCo (announced 2009-04-20, cash-or-stock election) — raw filings to a priced trade")
    steps = [
        ("1  Deal input", "Bloomberg row",
         "PepsiAmericas Inc  ←  PepsiCo Inc  ·  announced 2009-04-20  ·  Cash-or-Stock election  ·  status: completed"),
        ("2  SEC retrieval", "download_ma_edgar_files.py",
         "Resolved CIK → pulled 553 documents / 104 unique filings (8-K, 425, DEFA14A, SC 13E3/A) → top 10 scored & sent to Claude"),
        ("3  Claude extraction", "17 canonical fields",
         "Cash $28.50/sh  ·  Stock 0.5022 PEP sh (fixed ratio)  ·  Cash cap 50% of shares  ·  deadline = 3rd business day pre-close"),
        ("4  Realized label", "normalized_labels.csv",
         "Filing discloses DEMAND: ≈ 2,304,733 PAS shares elected cash → % elected cash (pre-proration) — the model's target variable"),
        ("5  WRDS join", "wrds_market_daily.csv",
         "Target + acquirer (PepsiCo) daily prices joined — 85 trading days each; supplies the entry price and the hedge leg"),
        ("6  Trade decision", "arb_signal.py",
         "Fair value $27.99 vs market $24.66 = +13.5% → risk-adjusted +9.1% → ENTER, elect CASH, hedge 0.066  ·  realized +15.6%"),
    ]
    y = 1.66
    for label, sub, detail in steps:
        rect(s, 0.70, y, 2.75, 0.80, BOXBG)
        txt(s, label, 0.80, y + 0.14, 2.55, 0.20, 10.5, NAVY, bold=True)
        txt(s, sub, 0.80, y + 0.42, 2.55, 0.20, 8.2, SLATE)
        rect(s, 3.62, y, 9.03, 0.80, WHITE, line=BORDER)
        txt(s, detail, 3.78, y + 0.17, 8.75, 0.50, 10.2, DARK)
        y += 0.895
    txt(s, "23 / 14", 12.10, 7.08, 0.80, 0.20, 7.8, SLATE, align=PP_ALIGN.RIGHT)
    return s


def reorder_before_last(prs, new_ids):
    """Move the two just-appended slides to sit right before the final (next-steps) slide."""
    lst = prs.slides._sldIdLst
    ids = list(lst)
    nextsteps = ids[-3]  # before appending, last was next-steps; after append it's 3rd from end
    a, b = ids[-2], ids[-1]
    lst.remove(a); lst.remove(b)
    pos = list(lst).index(nextsteps)
    lst.insert(pos, a); lst.insert(pos + 1, b)


def renumber(prs):
    total = len(prs.slides) - 1
    for idx, slide in enumerate(prs.slides):
        for sh in slide.shapes:
            if sh.has_text_frame:
                t = sh.text_frame.text.strip()
                if re.fullmatch(r"\d+\s*/\s*\d+", t):
                    p = sh.text_frame.paragraphs[0]
                    if p.runs:
                        p.runs[0].text = f"{idx} / {total}"
                        for r in p.runs[1:]:
                            r.text = ""


def main():
    prs = Presentation(SRC)
    update_weekly(prs)
    build_pipeline(prs)
    build_walkthrough(prs)
    reorder_before_last(prs, None)
    renumber(prs)
    prs.save(SRC)
    print(f"[deck] saved {SRC} — {len(prs.slides)} slides")
    for i, s in enumerate(prs.slides):
        title = next((sh.text_frame.text.split("\n")[0][:52] for sh in s.shapes
                      if sh.has_text_frame and sh.text_frame.text.strip()), "")
        print(f"  {i:2d}  {title}")


if __name__ == "__main__":
    main()
